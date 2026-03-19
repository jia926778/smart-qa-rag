"""多格式文档加载器模块，支持高级内容提取能力。

支持的文件格式和提取策略：

| 类别        | 扩展名                        | 策略 / 依赖库                              |
|------------|------------------------------|-------------------------------------------|
| PDF        | .pdf                         | pdfplumber (表格+文本), 降级 PyPDF         |
| Word       | .docx                        | docx2txt / python-docx                    |
| Excel      | .xlsx, .xls, .csv            | openpyxl / pandas → 文本表格               |
| PowerPoint | .pptx                        | python-pptx 幻灯片文本 + 备注              |
| Markdown   | .md                          | UnstructuredMarkdownLoader                |
| 纯文本      | .txt, .log, .json, .xml      | TextLoader (UTF-8)                        |
| HTML / 网页 | .html, .htm                  | BeautifulSoup 文本提取                     |
| 图片/OCR   | .png, .jpg, .jpeg, .tiff, .bmp | pytesseract OCR                          |
| 音视频      | .mp3, .wav, .mp4, .m4a, .webm  | OpenAI Whisper 转写                       |

每个加载器返回 ``List[Document]``，包含丰富的元数据（source, page,
file_type, extraction_method）。所有外部依赖均为延迟导入，
如果可选包未安装，系统会优雅降级。

主要组件：
- DocumentLoaderFactory: 文档加载器工厂类，根据文件扩展名选择合适的加载器
"""

from __future__ import annotations

import csv
import io
import json
import os
import tempfile
from pathlib import Path
from typing import Any, Dict, List, Optional

from langchain_core.documents import Document

from app.utils.exceptions import DocumentLoadError, UnsupportedFileTypeError
from app.utils.logger import get_logger

logger = get_logger(__name__)


# ---------------------------------------------------------------------------
# 基础辅助函数
# ---------------------------------------------------------------------------

def _meta(file_path: str, **extra: Any) -> Dict[str, Any]:
    """构建文档的基础元数据字典。
    
    Args:
        file_path: 文件路径。
        **extra: 额外的元数据字段。
    
    Returns:
        包含 source、file_type 及额外字段的元数据字典。
    """
    p = Path(file_path)
    return {
        "source": p.name,
        "file_type": p.suffix.lower().lstrip("."),
        **extra,
    }


# ---------------------------------------------------------------------------
# PDF 加载器 — pdfplumber 支持表格提取，降级到 PyPDF
# ---------------------------------------------------------------------------

def _load_pdf(file_path: str) -> List[Document]:
    """加载 PDF 文件，提取文本和表格内容。
    
    优先使用 pdfplumber（支持表格提取），失败时降级到 PyPDF。
    
    Args:
        file_path: PDF 文件路径。
    
    Returns:
        文档列表，每页一个 Document 对象。
    
    Raises:
        DocumentLoadError: PDF 加载失败时抛出。
    """
    docs: List[Document] = []
    try:
        import pdfplumber  # type: ignore

        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                parts: List[str] = []

                # 提取文本
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(text)

                # 提取表格并转换为 Markdown 格式
                tables = page.extract_tables() or []
                for table in tables:
                    if not table:
                        continue
                    md_table = _table_to_markdown(table)
                    if md_table:
                        parts.append(f"\n[表格]\n{md_table}")

                if parts:
                    content = "\n".join(parts)
                    docs.append(Document(
                        page_content=content,
                        metadata=_meta(file_path, page=page_num, extraction_method="pdfplumber"),
                    ))

        if docs:
            logger.info("PDF loaded via pdfplumber: %d pages from %s", len(docs), file_path)
            return docs
    except ImportError:
        logger.info("pdfplumber not available, falling back to PyPDF")
    except Exception as exc:
        logger.warning("pdfplumber failed for %s: %s. Falling back to PyPDF.", file_path, exc)

    # 降级方案：PyPDF
    try:
        from langchain_community.document_loaders import PyPDFLoader
        loader = PyPDFLoader(file_path)
        docs = loader.load()
        for doc in docs:
            doc.metadata.update(_meta(file_path, extraction_method="pypdf"))
        logger.info("PDF loaded via PyPDF: %d pages from %s", len(docs), file_path)
        return docs
    except Exception as exc:
        raise DocumentLoadError(f"PDF load failed: {exc}") from exc


def _table_to_markdown(table: List[List[Optional[str]]]) -> str:
    """将二维表格转换为 Markdown 格式。
    
    Args:
        table: 二维列表，每个元素为单元格内容。
    
    Returns:
        Markdown 格式的表格字符串。
    """
    if not table or len(table) < 1:
        return ""
    # 清理单元格内容
    cleaned = []
    for row in table:
        cleaned.append([str(cell).strip() if cell else "" for cell in row])

    lines = []
    # 表头
    lines.append("| " + " | ".join(cleaned[0]) + " |")
    lines.append("| " + " | ".join(["---"] * len(cleaned[0])) + " |")
    # 表体
    for row in cleaned[1:]:
        # 如果行长度不足，用空字符串填充
        while len(row) < len(cleaned[0]):
            row.append("")
        lines.append("| " + " | ".join(row[:len(cleaned[0])]) + " |")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Word 加载器 (.docx)
# ---------------------------------------------------------------------------

def _load_docx(file_path: str) -> List[Document]:
    """加载 Word 文档 (.docx)。
    
    优先使用 docx2txt，失败时降级到 python-docx。
    
    Args:
        file_path: Word 文件路径。
    
    Returns:
        文档列表。
    
    Raises:
        DocumentLoadError: 加载失败时抛出。
    """
    try:
        import docx2txt  # type: ignore
        text = docx2txt.process(file_path)
        if text and text.strip():
            return [Document(
                page_content=text,
                metadata=_meta(file_path, extraction_method="docx2txt"),
            )]
        return []
    except ImportError:
        pass

    # 降级方案：python-docx（提供更多结构信息）
    try:
        from docx import Document as DocxDocument  # type: ignore
        doc = DocxDocument(file_path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if paragraphs:
            return [Document(
                page_content="\n".join(paragraphs),
                metadata=_meta(file_path, extraction_method="python-docx"),
            )]
        return []
    except ImportError:
        raise DocumentLoadError("Neither docx2txt nor python-docx is installed")
    except Exception as exc:
        raise DocumentLoadError(f"DOCX load failed: {exc}") from exc


# ---------------------------------------------------------------------------
# Excel 加载器 (.xlsx, .xls, .csv)
# ---------------------------------------------------------------------------

def _load_excel(file_path: str) -> List[Document]:
    """加载 Excel 文件 (.xlsx, .xls) 或 CSV 文件。
    
    Args:
        file_path: Excel 或 CSV 文件路径。
    
    Returns:
        文档列表，每个工作表一个 Document 对象。
    
    Raises:
        DocumentLoadError: 加载失败时抛出。
    """
    ext = Path(file_path).suffix.lower()
    try:
        import openpyxl  # type: ignore
    except ImportError:
        openpyxl = None

    docs: List[Document] = []

    if ext == ".csv":
        return _load_csv(file_path)

    if ext in (".xlsx", ".xls"):
        # 优先使用 openpyxl
        if openpyxl:
            try:
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    rows = []
                    for row in ws.iter_rows(values_only=True):
                        rows.append([str(cell) if cell is not None else "" for cell in row])
                    if rows:
                        md = _table_to_markdown(rows)
                        if md:
                            docs.append(Document(
                                page_content=f"## Sheet: {sheet_name}\n\n{md}",
                                metadata=_meta(file_path, sheet=sheet_name, extraction_method="openpyxl"),
                            ))
                wb.close()
                if docs:
                    logger.info("Excel loaded via openpyxl: %d sheet(s) from %s", len(docs), file_path)
                    return docs
            except Exception as exc:
                logger.warning("openpyxl failed: %s", exc)

        # 降级方案：pandas
        try:
            import pandas as pd  # type: ignore
            xls = pd.ExcelFile(file_path)
            for sheet_name in xls.sheet_names:
                df = pd.read_excel(xls, sheet_name=sheet_name)
                text = f"## Sheet: {sheet_name}\n\n{df.to_markdown(index=False)}"
                docs.append(Document(
                    page_content=text,
                    metadata=_meta(file_path, sheet=sheet_name, extraction_method="pandas"),
                ))
            if docs:
                return docs
        except ImportError:
            raise DocumentLoadError("openpyxl or pandas required for Excel files")
        except Exception as exc:
            raise DocumentLoadError(f"Excel load failed: {exc}") from exc

    return docs


def _load_csv(file_path: str) -> List[Document]:
    """加载 CSV 文件。
    
    Args:
        file_path: CSV 文件路径。
    
    Returns:
        文档列表，包含表格的 Markdown 表示。
    
    Raises:
        DocumentLoadError: 加载失败时抛出。
    """
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            rows = [row for row in reader]
        if rows:
            md = _table_to_markdown(rows)
            return [Document(
                page_content=md,
                metadata=_meta(file_path, extraction_method="csv_reader"),
            )]
        return []
    except Exception as exc:
        raise DocumentLoadError(f"CSV load failed: {exc}") from exc


# ---------------------------------------------------------------------------
# PowerPoint 加载器 (.pptx)
# ---------------------------------------------------------------------------

def _load_pptx(file_path: str) -> List[Document]:
    """加载 PowerPoint 文件 (.pptx)。
    
    提取每张幻灯片的文本、表格和演讲者备注。
    
    Args:
        file_path: PowerPoint 文件路径。
    
    Returns:
        文档列表，每张幻灯片一个 Document 对象。
    
    Raises:
        DocumentLoadError: 加载失败时抛出。
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise DocumentLoadError("python-pptx is required for .pptx files")

    try:
        prs = Presentation(file_path)
        docs: List[Document] = []
        for slide_num, slide in enumerate(prs.slides, 1):
            parts: List[str] = []
            for shape in slide.shapes:
                # 提取文本框内容
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)
                # 提取表格
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        rows.append([cell.text.strip() for cell in row.cells])
                    md = _table_to_markdown(rows)
                    if md:
                        parts.append(md)
            # 提取演讲者备注
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[演讲者备注] {notes}")

            if parts:
                docs.append(Document(
                    page_content="\n".join(parts),
                    metadata=_meta(file_path, slide=slide_num, extraction_method="python-pptx"),
                ))
        logger.info("PPTX loaded: %d slide(s) from %s", len(docs), file_path)
        return docs
    except DocumentLoadError:
        raise
    except Exception as exc:
        raise DocumentLoadError(f"PPTX load failed: {exc}") from exc


# ---------------------------------------------------------------------------
# Markdown 加载器
# ---------------------------------------------------------------------------

def _load_markdown(file_path: str) -> List[Document]:
    """加载 Markdown 文件。
    
    Args:
        file_path: Markdown 文件路径。
    
    Returns:
        文档列表。
    
    Raises:
        DocumentLoadError: 加载失败时抛出。
    """
    try:
        from langchain_community.document_loaders import UnstructuredMarkdownLoader
        loader = UnstructuredMarkdownLoader(file_path)
        docs = loader.load()
        for doc in docs:
            doc.metadata.update(_meta(file_path, extraction_method="unstructured_md"))
        return docs
    except ImportError:
        # 降级方案：作为纯文本加载
        return _load_text(file_path)
    except Exception as exc:
        raise DocumentLoadError(f"Markdown load failed: {exc}") from exc


# ---------------------------------------------------------------------------
# 纯文本 / 结构化文本加载器
# ---------------------------------------------------------------------------

def _load_text(file_path: str) -> List[Document]:
    """加载纯文本文件。
    
    Args:
        file_path: 文本文件路径。
    
    Returns:
        文档列表。
    
    Raises:
        DocumentLoadError: 加载失败时抛出。
    """
    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        if not text.strip():
            return []
        return [Document(
            page_content=text,
            metadata=_meta(file_path, extraction_method="text_reader"),
        )]
    except Exception as exc:
        raise DocumentLoadError(f"Text load failed: {exc}") from exc


def _load_json(file_path: str) -> List[Document]:
    """加载 JSON 文件并转换为可读文本。
    
    Args:
        file_path: JSON 文件路径。
    
    Returns:
        文档列表，包含格式化的 JSON 文本。
    
    Raises:
        DocumentLoadError: 加载失败时抛出。
    """
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            data = json.load(f)
        text = json.dumps(data, ensure_ascii=False, indent=2)
        return [Document(
            page_content=text,
            metadata=_meta(file_path, extraction_method="json_reader"),
        )]
    except Exception as exc:
        raise DocumentLoadError(f"JSON load failed: {exc}") from exc


# ---------------------------------------------------------------------------
# HTML / 网页加载器
# ---------------------------------------------------------------------------

def _load_html(file_path: str) -> List[Document]:
    """加载 HTML 文件，提取文本和表格内容。
    
    Args:
        file_path: HTML 文件路径。
    
    Returns:
        文档列表。
    
    Raises:
        DocumentLoadError: 加载失败时抛出。
    """
    try:
        from bs4 import BeautifulSoup  # type: ignore
    except ImportError:
        # 降级方案：使用正则表达式去除标签
        return _load_html_fallback(file_path)

    try:
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            html = f.read()
        soup = BeautifulSoup(html, "html.parser")

        # 移除脚本、样式、导航等非内容元素
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()

        # 单独提取表格
        parts: List[str] = []
        for table_tag in soup.find_all("table"):
            rows = []
            for tr in table_tag.find_all("tr"):
                cells = [td.get_text(strip=True) for td in tr.find_all(["td", "th"])]
                rows.append(cells)
            if rows:
                parts.append(_table_to_markdown(rows))
            table_tag.decompose()

        # 提取主要文本
        main_text = soup.get_text(separator="\n", strip=True)
        if main_text:
            parts.insert(0, main_text)

        content = "\n\n".join(parts)
        if content.strip():
            return [Document(
                page_content=content,
                metadata=_meta(file_path, extraction_method="beautifulsoup"),
            )]
        return []
    except Exception as exc:
        raise DocumentLoadError(f"HTML load failed: {exc}") from exc


def _load_html_fallback(file_path: str) -> List[Document]:
    """HTML 加载的降级方案，使用正则表达式去除标签。
    
    Args:
        file_path: HTML 文件路径。
    
    Returns:
        文档列表。
    """
    import re
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        html = f.read()
    text = re.sub(r"<[^>]+>", " ", html)
    text = re.sub(r"\s+", " ", text).strip()
    if text:
        return [Document(page_content=text, metadata=_meta(file_path, extraction_method="regex_strip"))]
    return []


# ---------------------------------------------------------------------------
# OCR 加载器（扫描图片 / PDF）
# ---------------------------------------------------------------------------

def _load_ocr_image(file_path: str) -> List[Document]:
    """使用 OCR 技术从图片中提取文本。
    
    Args:
        file_path: 图片文件路径。
    
    Returns:
        文档列表。
    
    Raises:
        DocumentLoadError: OCR 失败时抛出。
    """
    try:
        from PIL import Image  # type: ignore
        import pytesseract  # type: ignore
    except ImportError:
        raise DocumentLoadError(
            "Pillow and pytesseract are required for OCR. "
            "Install: pip install Pillow pytesseract"
        )

    try:
        image = Image.open(file_path)
        # 尝试中文 + 英文 OCR
        text = pytesseract.image_to_string(image, lang="chi_sim+eng")
        if text and text.strip():
            return [Document(
                page_content=text.strip(),
                metadata=_meta(file_path, extraction_method="pytesseract_ocr"),
            )]
        return []
    except Exception as exc:
        raise DocumentLoadError(f"OCR failed: {exc}") from exc


def _load_scanned_pdf(file_path: str) -> List[Document]:
    """对扫描版 PDF 进行 OCR 处理。
    
    将 PDF 页面转换为图片后进行 OCR 识别。
    
    Args:
        file_path: PDF 文件路径。
    
    Returns:
        文档列表，每页一个 Document 对象。
    
    Raises:
        DocumentLoadError: OCR 失败时抛出。
    """
    try:
        from pdf2image import convert_from_path  # type: ignore
        import pytesseract  # type: ignore
    except ImportError:
        raise DocumentLoadError(
            "pdf2image and pytesseract are required for scanned PDF OCR. "
            "Install: pip install pdf2image pytesseract"
        )

    try:
        images = convert_from_path(file_path, dpi=300)
        docs: List[Document] = []
        for page_num, img in enumerate(images, 1):
            text = pytesseract.image_to_string(img, lang="chi_sim+eng")
            if text and text.strip():
                docs.append(Document(
                    page_content=text.strip(),
                    metadata=_meta(file_path, page=page_num, extraction_method="pdf_ocr"),
                ))
        logger.info("Scanned PDF OCR: %d page(s) from %s", len(docs), file_path)
        return docs
    except Exception as exc:
        raise DocumentLoadError(f"Scanned PDF OCR failed: {exc}") from exc


# ---------------------------------------------------------------------------
# 音视频转写（Whisper）
# ---------------------------------------------------------------------------

def _load_audio_video(file_path: str, whisper_model: str = "base") -> List[Document]:
    """使用 OpenAI Whisper 对音视频文件进行转写。
    
    Args:
        file_path: 音视频文件路径。
        whisper_model: Whisper 模型名称，默认为 "base"。
    
    Returns:
        文档列表，长音频按时间窗口分段。
    
    Raises:
        DocumentLoadError: 转写失败时抛出。
    """
    try:
        import whisper  # type: ignore
    except ImportError:
        raise DocumentLoadError(
            "openai-whisper is required for audio/video transcription. "
            "Install: pip install openai-whisper"
        )

    try:
        logger.info("Transcribing with Whisper (model=%s): %s", whisper_model, file_path)
        model = whisper.load_model(whisper_model)
        result = model.transcribe(file_path, language=None)  # 自动检测语言

        text = result.get("text", "")
        language = result.get("language", "unknown")
        segments = result.get("segments", [])

        docs: List[Document] = []

        if not text or not text.strip():
            return []

        # 如果音频较长，按片段分段以便更好地切分
        if len(segments) > 10:
            # 将片段按约 60 秒的时间窗口分组
            current_parts: List[str] = []
            current_start = 0.0
            window_size = 60.0  # 秒

            for seg in segments:
                seg_text = seg.get("text", "").strip()
                seg_start = seg.get("start", 0.0)
                if not seg_text:
                    continue

                if seg_start - current_start > window_size and current_parts:
                    docs.append(Document(
                        page_content=" ".join(current_parts),
                        metadata=_meta(
                            file_path,
                            start_time=current_start,
                            end_time=seg_start,
                            language=language,
                            extraction_method="whisper",
                        ),
                    ))
                    current_parts = []
                    current_start = seg_start

                current_parts.append(seg_text)

            # 处理剩余部分
            if current_parts:
                docs.append(Document(
                    page_content=" ".join(current_parts),
                    metadata=_meta(
                        file_path,
                        start_time=current_start,
                        language=language,
                        extraction_method="whisper",
                    ),
                ))
        else:
            docs.append(Document(
                page_content=text.strip(),
                metadata=_meta(file_path, language=language, extraction_method="whisper"),
            ))

        logger.info(
            "Whisper transcription: %d segment(s), lang=%s from %s",
            len(docs), language, file_path,
        )
        return docs
    except DocumentLoadError:
        raise
    except Exception as exc:
        raise DocumentLoadError(f"Audio/video transcription failed: {exc}") from exc


# ---------------------------------------------------------------------------
# Unstructured 通用降级加载器
# ---------------------------------------------------------------------------

def _load_unstructured(file_path: str) -> List[Document]:
    """使用 unstructured 库的自动检测功能作为最后的降级方案。
    
    Args:
        file_path: 文件路径。
    
    Returns:
        文档列表。
    
    Raises:
        DocumentLoadError: 加载失败时抛出。
    """
    try:
        from unstructured.partition.auto import partition  # type: ignore
    except ImportError:
        raise DocumentLoadError(
            "unstructured library is required as fallback. "
            "Install: pip install unstructured"
        )

    try:
        elements = partition(filename=file_path)
        text = "\n\n".join(str(el) for el in elements if str(el).strip())
        if text.strip():
            return [Document(
                page_content=text,
                metadata=_meta(file_path, extraction_method="unstructured_auto"),
            )]
        return []
    except Exception as exc:
        raise DocumentLoadError(f"Unstructured fallback failed: {exc}") from exc


# ===================================================================
# 主工厂类
# ===================================================================

# 扩展名 -> 加载函数映射
_LOADER_REGISTRY: Dict[str, Any] = {
    # PDF
    ".pdf": _load_pdf,
    # Word
    ".docx": _load_docx,
    # Excel / CSV
    ".xlsx": _load_excel,
    ".xls": _load_excel,
    ".csv": _load_excel,
    # PowerPoint
    ".pptx": _load_pptx,
    # Markdown
    ".md": _load_markdown,
    # 纯文本 / 结构化
    ".txt": _load_text,
    ".log": _load_text,
    ".json": _load_json,
    ".xml": _load_text,
    ".yaml": _load_text,
    ".yml": _load_text,
    # HTML
    ".html": _load_html,
    ".htm": _load_html,
    # 图片 (OCR)
    ".png": _load_ocr_image,
    ".jpg": _load_ocr_image,
    ".jpeg": _load_ocr_image,
    ".tiff": _load_ocr_image,
    ".tif": _load_ocr_image,
    ".bmp": _load_ocr_image,
    # 音视频 (Whisper 转写)
    ".mp3": _load_audio_video,
    ".wav": _load_audio_video,
    ".m4a": _load_audio_video,
    ".mp4": _load_audio_video,
    ".webm": _load_audio_video,
    ".flac": _load_audio_video,
    ".ogg": _load_audio_video,
}


class DocumentLoaderFactory:
    """文档加载器工厂类，根据文件扩展名选择合适的加载器。
    
    支持 25+ 种文件扩展名，涵盖 8 大类别：
    PDF、Word、Excel/CSV、PowerPoint、Markdown、HTML、图片(OCR)、
    音视频(Whisper 转写)。
    """

    @staticmethod
    def supported_extensions() -> List[str]:
        """获取所有支持的文件扩展名列表。
        
        Returns:
            排序后的扩展名列表。
        """
        return sorted(_LOADER_REGISTRY.keys())

    @staticmethod
    def supported_categories() -> Dict[str, List[str]]:
        """获取按类别分组的支持扩展名。
        
        Returns:
            类别名到扩展名列表的映射字典。
        """
        return {
            "PDF (含表格提取)": [".pdf"],
            "Word 文档": [".docx"],
            "Excel / CSV": [".xlsx", ".xls", ".csv"],
            "PowerPoint": [".pptx"],
            "Markdown": [".md"],
            "纯文本 / 结构化": [".txt", ".log", ".json", ".xml", ".yaml", ".yml"],
            "HTML 网页": [".html", ".htm"],
            "图片 (OCR)": [".png", ".jpg", ".jpeg", ".tiff", ".tif", ".bmp"],
            "音视频 (Whisper 转写)": [".mp3", ".wav", ".m4a", ".mp4", ".webm", ".flac", ".ogg"],
        }

    @staticmethod
    def load(file_path: str, enable_ocr_fallback: bool = True) -> List[Document]:
        """根据文件扩展名加载文档。
        
        Args:
            file_path: 文件路径。
            enable_ocr_fallback: 是否启用 PDF OCR 降级，默认为 True。
                当 PDF 正常提取无文本时，尝试 OCR 方式。
        
        Returns:
            文档列表，包含文本内容和元数据。
        
        Raises:
            UnsupportedFileTypeError: 文件类型不支持时抛出。
            DocumentLoadError: 文档加载失败时抛出。
        """
        ext = Path(file_path).suffix.lower()
        loader_fn = _LOADER_REGISTRY.get(ext)

        if loader_fn is None:
            # 尝试使用 unstructured 作为通用降级方案
            try:
                logger.info("Unknown extension %s, trying unstructured fallback", ext)
                docs = _load_unstructured(file_path)
                if docs:
                    return docs
            except Exception:
                pass
            raise UnsupportedFileTypeError(ext)

        try:
            docs = loader_fn(file_path)

            # PDF OCR 降级：如果 PDF 未提取到文本，尝试扫描版 PDF OCR
            if ext == ".pdf" and enable_ocr_fallback and not docs:
                logger.info("PDF yielded no text, attempting OCR fallback for %s", file_path)
                try:
                    docs = _load_scanned_pdf(file_path)
                except Exception as ocr_exc:
                    logger.warning("OCR fallback also failed: %s", ocr_exc)

            if not docs:
                logger.warning("No content extracted from %s", file_path)
                docs = [Document(
                    page_content="[文件内容为空或无法提取]",
                    metadata=_meta(file_path, extraction_method="empty"),
                )]

            logger.info(
                "Loaded %d document(s) from %s [%s]",
                len(docs),
                file_path,
                docs[0].metadata.get("extraction_method", "?") if docs else "?",
            )
            return docs

        except (UnsupportedFileTypeError, DocumentLoadError):
            raise
        except Exception as exc:
            logger.error("Failed to load %s: %s", file_path, exc)
            raise DocumentLoadError(f"Failed to load {file_path}: {exc}") from exc
